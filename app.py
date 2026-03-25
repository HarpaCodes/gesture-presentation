"""
Gesture Controlled Presentation System - Flask Backend
Run: python3 app.py
Open: http://localhost:8080
"""

import os
import subprocess
import shutil
from flask import Flask, request, jsonify, send_from_directory, render_template

app = Flask(__name__)

UPLOAD_FOLDER = 'static/uploads'
SLIDES_FOLDER = 'static/slides'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SLIDES_FOLDER, exist_ok=True)

state = {
    "current_slide": 0,
    "total_slides": 0,
    "locked": False,
    "slides": [],
    "presentation_name": ""
}

ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'png', 'jpg', 'jpeg'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


# ─── Conversion pipeline ──────────────────────────────────────────────────────

def convert_to_slides(filepath, filename):
    ext = filename.rsplit('.', 1)[1].lower()
    slides_dir = os.path.join(SLIDES_FOLDER, 'current')
    os.makedirs(slides_dir, exist_ok=True)

    # Clear old slides
    for f in os.listdir(slides_dir):
        try: os.remove(os.path.join(slides_dir, f))
        except: pass

    if ext in ('pptx', 'ppt'):
        # Step 1: PPTX → PDF using macOS's built-in sips/qlmanage or LibreOffice
        pdf_path = convert_pptx_to_pdf(filepath, filename)
        if pdf_path:
            result = convert_pdf_to_images(pdf_path, slides_dir)
            if result: return result

        # Step 2: Try python-pptx screenshot method (slide thumbnails)
        result = convert_pptx_thumbnails(filepath, slides_dir)
        if result: return result

        print("❌ All PPTX methods failed.")
        print("👉 FIX: Open your PPTX in PowerPoint → File → Export → PDF → upload that PDF instead")
        return generate_demo_slides(slides_dir)

    if ext == 'pdf':
        result = convert_pdf_to_images(filepath, slides_dir)
        if result: return result
        print("❌ PDF conversion failed — install poppler: brew install poppler")
        return generate_demo_slides(slides_dir)

    if ext in ('png', 'jpg', 'jpeg'):
        out = os.path.join(slides_dir, 'slide_001.png')
        shutil.copy(filepath, out)
        return ['slides/current/slide_001.png']

    return generate_demo_slides(slides_dir)


def convert_pptx_to_pdf(filepath, filename):
    """Convert PPTX to PDF using available tools on the system."""
    pdf_name = filename.rsplit('.', 1)[0] + '.pdf'
    pdf_path = os.path.join(UPLOAD_FOLDER, pdf_name)

    # Method 1: LibreOffice (if installed)
    for lo_cmd in ['libreoffice', 'soffice', '/Applications/LibreOffice.app/Contents/MacOS/soffice']:
        try:
            result = subprocess.run(
                [lo_cmd, '--headless', '--convert-to', 'pdf',
                 '--outdir', os.path.abspath(UPLOAD_FOLDER),
                 os.path.abspath(filepath)],
                capture_output=True, timeout=120
            )
            if os.path.exists(pdf_path):
                print(f"✓ PPTX→PDF via LibreOffice")
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    # Method 2: macOS qlmanage (QuickLook — generates thumbnails, limited)
    # Method 3: unoconv
    try:
        result = subprocess.run(
            ['unoconv', '-f', 'pdf', '-o', pdf_path, filepath],
            capture_output=True, timeout=120
        )
        if os.path.exists(pdf_path):
            print("✓ PPTX→PDF via unoconv")
            return pdf_path
    except FileNotFoundError:
        pass

    return None


def convert_pdf_to_images(pdf_path, slides_dir):
    """Convert PDF pages to PNG images."""

    # Method 1: pdf2image (requires poppler)
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path, dpi=150)
        paths = []
        for i, img in enumerate(images):
            out = os.path.join(slides_dir, f'slide_{i+1:03d}.png')
            img.save(out, 'PNG')
            paths.append(f'slides/current/slide_{i+1:03d}.png')
        print(f"✓ PDF→PNG via pdf2image: {len(paths)} slides")
        return paths
    except Exception as e:
        print(f"pdf2image failed: {e}")

    # Method 2: macOS built-in `sips` + `qlmanage`  
    try:
        result = subprocess.run(
            ['qlmanage', '-t', '-s', '1280', '-o', slides_dir, pdf_path],
            capture_output=True, timeout=60
        )
        # qlmanage outputs filename.pdf.png
        pngs = sorted([f for f in os.listdir(slides_dir) if f.endswith('.png')])
        if pngs:
            paths = []
            for i, f in enumerate(pngs):
                new_name = f'slide_{i+1:03d}.png'
                os.rename(os.path.join(slides_dir, f), os.path.join(slides_dir, new_name))
                paths.append(f'slides/current/{new_name}')
            print(f"✓ PDF→PNG via qlmanage: {len(paths)} slides")
            return paths
    except Exception as e:
        print(f"qlmanage failed: {e}")

    # Method 3: PyMuPDF (fitz) - no external tools needed!
    try:
        import fitz  # pip install pymupdf
        doc = fitz.open(pdf_path)
        paths = []
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2, 2)  # 2x zoom = ~150dpi
            pix = page.get_pixmap(matrix=mat)
            out = os.path.join(slides_dir, f'slide_{i+1:03d}.png')
            pix.save(out)
            paths.append(f'slides/current/slide_{i+1:03d}.png')
        doc.close()
        print(f"✓ PDF→PNG via PyMuPDF: {len(paths)} slides")
        return paths
    except ImportError:
        print("PyMuPDF not installed. Try: pip3 install pymupdf")
    except Exception as e:
        print(f"PyMuPDF failed: {e}")

    return None


def convert_pptx_thumbnails(filepath, slides_dir):
    """
    Use python-pptx to extract slide info + render basic thumbnails.
    Not pixel-perfect but shows real content structure.
    """
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        import io

        prs = Presentation(filepath)
        W, H = 1280, 720
        paths = []

        for idx, slide in enumerate(prs.slides):
            img = Image.new('RGB', (W, H), (15, 20, 35))
            draw = ImageDraw.Draw(img)

            # Try to get background color
            try:
                bg = slide.background.fill
                if bg.type is not None:
                    rgb = bg.fore_color.rgb
                    img = Image.new('RGB', (W, H), (rgb.r, rgb.g, rgb.b))
                    draw = ImageDraw.Draw(img)
            except: pass

            # Draw all text from the slide
            y_pos = 60
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text: continue

                    # Determine font size based on shape position/size
                    font_size = 28
                    try:
                        if shape.width > prs.slide_width * 0.6:
                            font_size = 48  # wide shapes = title
                        if para.runs and para.runs[0].font.size:
                            font_size = max(14, min(int(para.runs[0].font.size.pt * 1.3), 72))
                    except: pass

                    # Color
                    color = (220, 220, 220)
                    try:
                        rgb = para.runs[0].font.color.rgb
                        color = (rgb.r, rgb.g, rgb.b)
                    except: pass

                    # Try system fonts
                    font = None
                    for font_path in [
                        '/System/Library/Fonts/Helvetica.ttc',
                        '/System/Library/Fonts/Arial.ttf',
                        '/Library/Fonts/Arial.ttf',
                        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
                    ]:
                        try:
                            from PIL import ImageFont
                            font = ImageFont.truetype(font_path, font_size)
                            break
                        except: continue

                    if font:
                        draw.text((60, y_pos), text[:80], fill=color, font=font)
                        y_pos += font_size + 10
                    else:
                        draw.text((60, y_pos), text[:80], fill=color)
                        y_pos += 30

                    if y_pos > H - 40: break

            # Slide number
            draw.rectangle([0, H-36, W, H], fill=(20, 30, 50))
            try:
                from PIL import ImageFont
                fn = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 18)
                draw.text((W//2, H-18), f"Slide {idx+1} of {len(prs.slides)}",
                          fill=(100,150,200), font=fn, anchor='mm')
            except:
                draw.text((W//2 - 40, H-28), f"Slide {idx+1}", fill=(100,150,200))

            out = os.path.join(slides_dir, f'slide_{idx+1:03d}.png')
            img.save(out, 'PNG')
            paths.append(f'slides/current/slide_{idx+1:03d}.png')

        print(f"✓ PPTX thumbnails via python-pptx: {len(paths)} slides")
        return paths
    except Exception as e:
        print(f"python-pptx thumbnail failed: {e}")
        return None


def generate_demo_slides(slides_dir):
    """Colorful demo slides as last resort."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        data = [
            ("#1a1a2e","#e94560","Upload a PDF for best results","python3 app.py is running!"),
            ("#0f3460","#00d4ff","Tip: Export PPTX as PDF","File → Export → PDF in PowerPoint"),
            ("#533483","#ffffff","Then upload the PDF","Slides will look exactly right"),
        ]
        paths = []
        for i, (bg, acc, t1, t2) in enumerate(data):
            img = Image.new('RGB', (1280, 720), bg)
            draw = ImageDraw.Draw(img)
            draw.rectangle([0,0,1280,720], outline=acc, width=6)
            try:
                f1 = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 52)
                f2 = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 28)
            except:
                f1 = f2 = None
            if f1:
                draw.text((640,300), t1, fill=acc, font=f1, anchor='mm')
                draw.text((640,400), t2, fill='#cccccc', font=f2, anchor='mm')
            else:
                draw.text((200,300), t1, fill=acc)
                draw.text((200,400), t2, fill='#cccccc')
            out = os.path.join(slides_dir, f'slide_{i+1:03d}.png')
            img.save(out)
            paths.append(f'slides/current/slide_{i+1:03d}.png')
        return paths
    except:
        return []


# ─── Routes ──────────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/viewer')
def viewer():
    return render_template('viewer.html')

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    if not file.filename or not allowed_file(file.filename):
        return jsonify({"error": "Invalid file"}), 400

    filename = file.filename
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    file.save(filepath)

    print(f"\n📂 Uploaded: {filename}")
    slides = convert_to_slides(filepath, filename)

    state.update({
        "slides": slides,
        "total_slides": len(slides),
        "current_slide": 0,
        "locked": False,
        "presentation_name": filename
    })

    return jsonify({
        "success": True,
        "total_slides": len(slides),
        "slides": slides,
        "message": f"Loaded {len(slides)} slides"
    })


@app.route('/slides/current/<path:filename>')
def serve_slide(filename):
    return send_from_directory('static/slides/current', filename)


@app.route('/state')
def get_state():
    return jsonify(state)


@app.route('/gesture', methods=['POST'])
def handle_gesture():
    data = request.get_json(silent=True) or {}
    gesture = data.get("gesture", "")

    if state["locked"] and gesture != "fist":
        return jsonify({"status": "locked", "current_slide": state["current_slide"], "locked": True, "changed": False})

    changed = False
    if gesture == "next" and state["current_slide"] < state["total_slides"] - 1:
        state["current_slide"] += 1; changed = True
    elif gesture == "prev" and state["current_slide"] > 0:
        state["current_slide"] -= 1; changed = True
    elif gesture == "fist":
        state["locked"] = not state["locked"]; changed = True

    return jsonify({"status": "ok", "current_slide": state["current_slide"],
                    "locked": state["locked"], "changed": changed})


@app.route('/slide/next', methods=['POST'])
def next_slide():
    if not state["locked"] and state["current_slide"] < state["total_slides"] - 1:
        state["current_slide"] += 1
    return jsonify(state)

@app.route('/slide/prev', methods=['POST'])
def prev_slide():
    if not state["locked"] and state["current_slide"] > 0:
        state["current_slide"] -= 1
    return jsonify(state)

@app.route('/slide/goto', methods=['POST'])
def goto_slide():
    data = request.get_json(silent=True) or {}
    idx = data.get("index", 0)
    if 0 <= idx < state["total_slides"]:
        state["current_slide"] = idx
    return jsonify(state)

@app.route('/toggle_lock', methods=['POST'])
def toggle_lock():
    state["locked"] = not state["locked"]
    return jsonify(state)


if __name__ == '__main__':
    print("=" * 50)
    print("  GestureSlide — Running!")
    print("  Open: http://localhost:8080")
    print("=" * 50)
    app.run(debug=True, host='0.0.0.0', port=8080)
